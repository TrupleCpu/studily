from pptx import Presentation
from io import BytesIO
import pdfplumber
from docx import Document

def extract_text_from_pdf(pdf_bytes):
    text = ""
    with pdfplumber.open(BytesIO(pdf_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text


def extract_text_from_ppt(pptx_bytes):
    try:
        prs = Presentation(BytesIO(pptx_bytes))
    except Exception as e:
        raise RuntimeError(f"Error opening Powerpoint file: {e}")
    
    extracted_data = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_content = {
            "slide_number": slide_num,
            "title": "",
            "text": [],
            "notes": ""
        }

        if slide.shapes.title and slide.shapes.title.text:
            slide_content["title"] = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content["text"].append(shape.text.strip())

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_content["notes"] = notes_text
            
        extracted_data.append(slide_content)
    
    return extracted_data

def extract_text_from_docx(docx_bytes):
    text = ""
    doc = Document(BytesIO(docx_bytes))

    for para in doc.paragraphs:
        text += para.text + "\n"

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            text += row_text + "\n"

    return text.strip()



